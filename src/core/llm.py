import asyncio
import os
from google import genai
from google.genai import types
from streamlit.runtime.uploaded_file_manager import UploadedFile

import PyPDF2
import docx
from pptx import Presentation


class LLM:
    """
    LLM wrapper for Gemini API calls.
    """
    def __init__(self, api_key: str, model_name: str = "gemini-2.0-flash"):
        self.model_name = model_name
        self.client = genai.Client(api_key=api_key)

    def generate_insight(self, prompt: str) -> str:
        """
        This function returns the response from the gemini sdk
        """
        response = self.client.models.generate_content(
            model=self.model_name,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.7,
                top_p=0.9,
                max_output_tokens=3000,
                tools=[{"google_search": {}}],

            )
        )
        return response.text

    async def generate_insight_async(self, prompt: str) -> str:
        """
        Async wrapper of the generate_insights function
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self.generate_insight, prompt)

    @staticmethod
    def extract_text_from_file(uploaded_file: UploadedFile) -> str:
        """
        A function for extraction of content from the uploaded file
        """
        ext = os.path.splitext(uploaded_file.name)[1].lower()

        if ext == '.txt':
            uploaded_file.seek(0)
            data = uploaded_file.read()
            return data.decode('utf-8')

        elif ext == '.pdf':
            text = ""
            uploaded_file.seek(0)
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text

        elif ext in ['.docx', '.doc']:
            if ext == '.doc':
                raise ValueError("Extraction for .doc files is not supported. Please convert to .docx format.")
            uploaded_file.seek(0)
            doc = docx.Document(uploaded_file)
            paragraphs = [para.text for para in doc.paragraphs if para.text]
            return "\n".join(paragraphs)

        elif ext == '.pptx':
            uploaded_file.seek(0)
            prs = Presentation(uploaded_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            raise ValueError(f"Unsupported file extension: {ext}")
